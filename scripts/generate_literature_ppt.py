#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SatelliteSimJulia 文献调研汇报 PPT 生成脚本

读取 _data.json, 生成 ~25 页全景调研汇报 PPT:
  - 封面 / 目录
  - 项目背景 + 调研目的
  - 数据源与方法论 + 全景统计图
  - 10 板块各 2 页(概览 + Top 论文)
  - 研究热点 / 空白矩阵
  - 路线图 + 总结
"""

import json
import os
import sys
from io import BytesIO

import matplotlib
matplotlib.use("Agg")  # 无 GUI 后端
import matplotlib.pyplot as plt
from matplotlib import font_manager

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# 路径与配置
# ---------------------------------------------------------------------------
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
OUT_DIR = os.path.join(PROJECT_DIR, "docs", "literature")
JSON_IN = os.path.join(OUT_DIR, "_data.json")
PPT_OUT = os.path.join(OUT_DIR, "SatelliteSimJulia文献调研汇报.pptx")

# 中文字体配置 (macOS)
def setup_chinese_font():
    """尝试设置 matplotlib 中文字体。"""
    candidates = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/Library/Fonts/Songti.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                font_manager.fontManager.addfont(path)
                name = font_manager.FontProperties(fname=path).get_name()
                plt.rcParams["font.sans-serif"] = [name, "Arial Unicode MS"]
                plt.rcParams["axes.unicode_minus"] = False
                return name
            except Exception:
                pass
    plt.rcParams["axes.unicode_minus"] = False
    return None


CN_FONT = setup_chinese_font()

# 配色方案(深空蓝主题,契合卫星主题)
COLOR_BG = RGBColor(0x0F, 0x1B, 0x2D)         # 深空蓝背景
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)   # 浅背景
COLOR_PRIMARY = RGBColor(0x00, 0x7A, 0xCC)    # 主色(科技蓝)
COLOR_ACCENT = RGBColor(0x00, 0xC8, 0xB4)     # 强调(青绿)
COLOR_WARN = RGBColor(0xFF, 0x8C, 0x42)       # 警示(橙)
COLOR_DANGER = RGBColor(0xE7, 0x4C, 0x3C)     # 危险(红)
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)       # 正文深灰
COLOR_TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF) # 浅色文字
COLOR_MUTED = RGBColor(0x7F, 0x8C, 0x8D)      # 次要灰

# 板块主色(每板块一个色,用于图表)
SECTION_COLORS = [
    "#007ACC", "#00C8B4", "#FF8C42", "#9B59B6", "#E74C3C",
    "#F1C40F", "#2ECC71", "#3498DB", "#E91E63", "#95A5A6",
]


# ---------------------------------------------------------------------------
# 图表生成
# ---------------------------------------------------------------------------
def make_section_count_chart(secs_meta):
    """各板块论文数柱状图。"""
    fig, ax = plt.subplots(figsize=(10, 5))
    names = [f"{s['id']}\n{s['name']}" for s in secs_meta]
    counts = [s["total"] for s in secs_meta]
    colors = SECTION_COLORS[:len(names)]
    bars = ax.bar(names, counts, color=colors, edgecolor="white", linewidth=1.2)
    for bar, cnt in zip(bars, counts):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 8,
                str(cnt), ha="center", va="bottom", fontsize=11, fontweight="bold")
    ax.set_ylabel("论文数", fontsize=12)
    ax.set_title("各技术板块论文数量分布", fontsize=14, fontweight="bold", pad=12)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_ylim(0, max(counts) * 1.15)
    plt.xticks(fontsize=9)
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def make_tier_stacked_chart(secs_meta):
    """各板块 tier 分布堆叠柱状图。"""
    fig, ax = plt.subplots(figsize=(10, 5))
    names = [s["id"] for s in secs_meta]
    t1 = [s["tier1"] for s in secs_meta]
    t2 = [s["tier2"] for s in secs_meta]
    t3 = [s["tier3"] for s in secs_meta]
    x = range(len(names))
    ax.bar(x, t1, label="★核心", color="#007ACC", edgecolor="white")
    ax.bar(x, t2, bottom=t1, label="☆相关", color="#00C8B4", edgecolor="white")
    bottom2 = [a + b for a, b in zip(t1, t2)]
    ax.bar(x, t3, bottom=bottom2, label="○借鉴", color="#FF8C42", edgecolor="white")
    ax.set_xticks(list(x))
    ax.set_xticklabels(names, fontsize=11)
    ax.set_ylabel("论文数", fontsize=12)
    ax.set_title("各板块相关性分层(★核心 / ☆相关 / ○借鉴)",
                 fontsize=14, fontweight="bold", pad=12)
    ax.legend(loc="upper right", fontsize=10)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def make_year_trend_chart(data):
    """全库年份分布趋势图。"""
    yd = data["meta"]["year_distribution"]
    years = sorted([int(y) for y in yd.keys() if y.isdigit() and 2005 <= int(y) <= 2026])
    counts = [yd.get(str(y), 0) for y in years]
    fig, ax = plt.subplots(figsize=(10, 4.5))
    ax.fill_between(years, counts, color="#007ACC", alpha=0.25)
    ax.plot(years, counts, color="#007ACC", linewidth=2.5, marker="o",
            markersize=5, markerfacecolor="#00C8B4")
    ax.set_xlabel("年份", fontsize=12)
    ax.set_ylabel("论文数", fontsize=12)
    ax.set_title("全库论文年份分布趋势(2005-2026)",
                 fontsize=14, fontweight="bold", pad=12)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.grid(axis="y", alpha=0.3)
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def make_heatmap_chart(secs_meta):
    """研究热度横向条形图(按热度排序)。"""
    sorted_secs = sorted(secs_meta, key=lambda s: s["total"])
    fig, ax = plt.subplots(figsize=(10, 5))
    names = [f"{s['id']} {s['name']}" for s in sorted_secs]
    counts = [s["total"] for s in sorted_secs]
    colors = []
    for c in counts:
        if c >= 400:
            colors.append("#E74C3C")  # 红海
        elif c >= 150:
            colors.append("#F39C12")  # 活跃
        else:
            colors.append("#2ECC71")  # 蓝海
    bars = ax.barh(names, counts, color=colors, edgecolor="white", linewidth=1)
    for bar, cnt in zip(bars, counts):
        ax.text(bar.get_width() + 6, bar.get_y() + bar.get_height() / 2,
                str(cnt), va="center", fontsize=10, fontweight="bold")
    ax.set_xlabel("论文数", fontsize=12)
    ax.set_title("研究热度排序(红:饱和 → 绿:蓝海)",
                 fontsize=14, fontweight="bold", pad=12)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.set_xlim(0, max(counts) * 1.15)
    plt.tight_layout()
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# PPT 幻灯片构建辅助
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_slide(prs, layout_idx=6):  # 6 = blank
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=COLOR_TEXT,
                align=PP_ALIGN.LEFT, font_name=None, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(font_size)
    f.bold = bold
    f.color.rgb = color
    if font_name:
        f.name = font_name
    return txBox


def add_bullets(slide, left, top, width, height, bullets, *,
                font_size=16, color=COLOR_TEXT, font_name=None,
                line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        # 支持 (text, level) 元组
        if isinstance(b, tuple):
            text, level = b
        else:
            text, level = b, 0
        prefix = "    " * level + ("• " if level == 0 else "– ")
        run.text = prefix + text
        f = run.font
        f.size = Pt(font_size - level * 2)
        f.color.rgb = color if level == 0 else COLOR_MUTED
        if font_name:
            f.name = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_title_bar(slide, title, subtitle=None):
    """标准页眉:左侧色块 + 标题。"""
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(7.5), COLOR_PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                title, font_size=28, bold=True, color=COLOR_TEXT)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.95), Inches(12), Inches(0.4),
                    subtitle, font_size=14, color=COLOR_MUTED)
    # 底部页脚
    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
                "SatelliteSimJulia 文献调研汇报  |  2026-07-03",
                font_size=9, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# 各幻灯片生成
# ---------------------------------------------------------------------------
def slide_cover(prs, data):
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG)
    # 装饰色块
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), COLOR_PRIMARY)
    add_rect(slide, Inches(0), Inches(7.35), SLIDE_W, Inches(0.15), COLOR_ACCENT)
    # 左侧装饰
    add_rect(slide, Inches(0), Inches(2.5), Inches(0.25), Inches(2.5), COLOR_ACCENT)
    # 主标题
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2),
                "SatelliteSimJulia", font_size=54, bold=True,
                color=COLOR_TEXT_LIGHT)
    add_textbox(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.8),
                "文献调研全景汇报", font_size=36, bold=True,
                color=COLOR_ACCENT)
    # 副信息
    total = data["meta"]["total_papers_in_db"]
    matched = data["meta"]["total_unique_matched"]
    n_secs = len(data["sections_meta"])
    add_textbox(
        slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.6),
        f"从 {total:,} 篇论文库中筛选出 {matched:,} 篇相关文献  ·  按 {n_secs} 个技术层级分类",
        font_size=18, color=COLOR_MUTED)
    add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
                "LEO 卫星星座仿真 + 可微优化 + AI 适配",
                font_size=14, color=COLOR_MUTED)
    add_textbox(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5),
                "汇报日期:2026-07-03",
                font_size=12, color=COLOR_MUTED)


def slide_toc(prs):
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, "目录", "本次汇报结构")
    toc = [
        ("01", "项目背景与调研目的", "为什么做这次调研"),
        ("02", "数据源与方法论", "15800 篇 → 10 板块的筛选逻辑"),
        ("03", "全景统计概览", "论文分布与年份趋势"),
        ("04", "十大板块文献全景", "每板块核心问题 + 代表论文"),
        ("05", "研究热点与空白矩阵", "红海 / 活跃 / 蓝海识别"),
        ("06", "三大研究空白与本项目创新", "可微优化 / PINN / LLM 编排"),
        ("07", "路线图与总结", "分阶段推进建议"),
    ]
    y = 1.7
    for idx, title, desc in toc:
        # 序号圆角块
        add_rect(slide, Inches(0.8), Inches(y), Inches(0.6), Inches(0.6),
                 COLOR_PRIMARY)
        add_textbox(slide, Inches(0.8), Inches(y), Inches(0.6), Inches(0.6),
                    idx, font_size=18, bold=True, color=COLOR_TEXT_LIGHT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.6), Inches(y - 0.05), Inches(10), Inches(0.4),
                    title, font_size=18, bold=True, color=COLOR_TEXT)
        add_textbox(slide, Inches(1.6), Inches(y + 0.3), Inches(10), Inches(0.35),
                    desc, font_size=12, color=COLOR_MUTED)
        y += 0.72


def slide_background(prs):
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, "01  项目背景与调研目的",
                  "SatelliteSimJulia 是什么,为什么要做这次文献调研")
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
                "▎项目定位", font_size=20, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, Inches(0.6), Inches(1.95), Inches(12.2), Inches(1.6), [
        "LEO 卫星星座仿真 + 可微优化 + AI 适配 的端到端流水线",
        "全链路用裸 Array{Float64,3} 衔接,用多重分派扩展,并可微分做梯度优化",
        "从 Walker 星座生成 → 轨道传播 → ISL/GSL 链路 → 拓扑/路由 → 流量/容量 → 指标",
    ], font_size=15)
    add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                "▎调研目的", font_size=20, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, Inches(0.6), Inches(4.25), Inches(12.2), Inches(2.0), [
        "摸清各技术层级的文献基线与对标标杆(如 Hypatia 117ms RTT)",
        "识别研究空白,定位本项目的创新点与论文潜力",
        "为后续开发提供可验证的数值基准与可复现实验模板",
        "按 10 个技术层级组织文献,与 src/ 模块一一对应,便于开发参考",
    ], font_size=15)


def slide_methodology(prs, data):
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, "02  数据源与方法论",
                  "15801 篇论文库 → 3020 篇命中 → 10 板块分类")
    total = data["meta"]["total_papers_in_db"]
    matched = data["meta"]["total_unique_matched"]
    # 左侧:数据源
    add_textbox(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(0.5),
                "▎数据源", font_size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, Inches(0.6), Inches(1.95), Inches(5.8), Inches(2.5), [
        f"本地论文库:共 {total:,} 篇",
        "字段:标题/年份/来源/arXiv 分类标签/语义聚类",
        "来源:arXiv / CCF / CAS / 航天 / 国际供 等 7 类",
        "时间跨度:1986 - 2026(近年为主)",
    ], font_size=13)
    # 右侧:筛选方法
    add_textbox(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(0.5),
                "▎筛选方法论", font_size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, Inches(6.9), Inches(1.95), Inches(5.8), Inches(2.5), [
        f"标题关键词 + 分类标签交叉过滤 → 命中 {matched:,} 篇",
        "三级相关性:★核心(标题强匹配)",
        "　　　　　　☆相关(标签+卫星上下文)",
        "　　　　　　○借鉴(相邻领域,限2021+)",
        "CSV 标题含逗号按列位置重建,避免错位",
    ], font_size=13)
    # 底部:板块划分原则
    add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0),
             RGBColor(0xEA, 0xF2, 0xFA))
    add_textbox(slide, Inches(0.7), Inches(4.95), Inches(12), Inches(0.5),
                "▎板块划分原则:按项目技术层级", font_size=16, bold=True,
                color=COLOR_PRIMARY)
    add_bullets(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(1.3), [
        "与 SatelliteSimJulia 的 src/ 模块一一对应(orbit / link / net / opt / lab)",
        "从底层物理(轨道)到顶层智能(LLM 编排),覆盖全栈",
        "每个板块产出:核心问题 → 对标基准 → 子主题分组 → 全量论文清单 → 研究空白",
    ], font_size=12)


def slide_stats_overview(prs, data, chart1, chart2):
    """全景统计 - 用两个图表。"""
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, "03  全景统计概览",
                  "各板块论文数 + 相关性分层")
    slide.shapes.add_picture(chart1, Inches(0.4), Inches(1.3),
                             width=Inches(6.3))
    slide.shapes.add_picture(chart2, Inches(6.9), Inches(1.3),
                             width=Inches(6.3))
    # 关键数字卡片
    total_all = sum(s["total"] for s in data["sections_meta"])
    add_rect(slide, Inches(0.5), Inches(6.0), Inches(4), Inches(0.9),
             COLOR_PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(6.05), Inches(4), Inches(0.45),
                f"{total_all}", font_size=28, bold=True,
                color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(4), Inches(0.35),
                "板块论文总数(含跨板块重复)",
                font_size=11, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4.7), Inches(6.0), Inches(4), Inches(0.9),
             COLOR_ACCENT)
    add_textbox(slide, Inches(4.7), Inches(6.05), Inches(4), Inches(0.45),
                f"{data['meta']['total_unique_matched']:,}",
                font_size=28, bold=True, color=COLOR_TEXT_LIGHT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.7), Inches(6.5), Inches(4), Inches(0.35),
                "去重后命中独立论文数",
                font_size=11, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(8.9), Inches(6.0), Inches(4), Inches(0.9),
             COLOR_WARN)
    add_textbox(slide, Inches(8.9), Inches(6.05), Inches(4), Inches(0.45),
                "10", font_size=28, bold=True, color=COLOR_TEXT_LIGHT,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.9), Inches(6.5), Inches(4), Inches(0.35),
                "技术层级板块",
                font_size=11, color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER)


def slide_year_trend(prs, data, chart):
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, "03  全景统计概览(续)",
                  "论文年份分布趋势")
    slide.shapes.add_picture(chart, Inches(1.5), Inches(1.4),
                             width=Inches(10.3))
    add_textbox(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
                "▎趋势观察", font_size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, Inches(0.6), Inches(6.3), Inches(12.2), Inches(1.0), [
        "2019 年后卫星网络论文爆发式增长,Starlink 商业化带动学术热度",
        "2024-2026 年达到峰值,反映 LEO 星座已成为网络研究最热方向",
        "近年(2021+)论文占比超 60%,说明本调研聚焦时效性合理",
    ], font_size=13)


# 板块概览的精简信息(每板块2页:概览 + Top论文)
SECTION_OVERVIEW = {
    "01": {
        "core": "如何在长时间尺度上高精度、低成本预测卫星位置?",
        "benchmarks": "Two-Body ~100km/天 · J2 ~10km/天(可微首选) · SGP4 ~3km · HPOP 亚米",
        "highlights": [
            "jaxsgp4: GPU 加速巨型星座批量传播",
            "NGSO Constellation Design for Global Connectivity",
            "Low-Earth Satellite Orbit Determination Using Deep CNN",
            "Minimum-hop Constellation Design for LEO",
        ],
        "gap": "可微 SGP4/J2 端到端梯度穿透到网络层优化几乎空白",
    },
    "02": {
        "core": "如何评估 ISL/GSL 链路物理可用性与质量?",
        "benchmarks": "Starlink 4 ISL/星(+Grid) · 光 ISL ~5ms/跳 · min_elevation=10°",
        "highlights": [
            "DeepISL: LEO 星间链路规划联合优化",
            "Investigating ISL Spanning Patterns on Networking Performance",
            "Enhancing Reliability via High-Speed Inter-Satellite Links",
            "On-Demand Routing with Dynamic Laser ISL",
        ],
        "gap": "动态激光 ISL 建立时延量化建模较少;可微链路评估空白",
    },
    "03": {
        "core": "如何设计 ISL 连接关系形成鲁棒低时延拓扑?",
        "benchmarks": "+Grid 4-ISL 标准 · 3-ISL 解析公式(南大) · Hypatia 100ms 重算",
        "highlights": [
            "Starfield: Demand-Aware Satellite Topology Design",
            "Time-Dependent Network Topology Optimization",
            "Topology Virtualization and Dynamics Shielding",
            "Exploiting Topology Awareness for Routing",
        ],
        "gap": "需求感知拓扑与可微拓扑联合优化是新兴方向",
    },
    "04": {
        "core": "如何在时变 LEO 拓扑上高效计算源-目路径?(最热门方向)",
        "benchmarks": "Hypatia RTT 85-117ms · PAM2023 跳数~10 · GraphSAGE +29%吞吐",
        "highlights": [
            "Hypatia / satgenpy(IMC 2020)— 首选对标",
            "A Fast Percolation-Dijkstra Routing for Mega-Constellation",
            "Segment Routing in Broadband LEO (Landmark Skeleton)",
            "Optimal Oblivious Load-Balancing for Sparse Traffic",
        ],
        "gap": "可微路由与 PINN 路由时延预测器(pinn_routing.jl)属首创",
    },
    "05": {
        "core": "如何评估端到端流量、链路负载、网络容量与时延分布?",
        "benchmarks": "Starlink 跨洲 RTT ~30-50ms · 光 ISL 10-100Gbps · Mathis 上界",
        "highlights": [
            "Optimal Oblivious Load-Balancing for Large-Scale Satellite",
            "Statistical Characterization of E2E Latency over LEO",
            "Throughput Analysis for Hybrid RF/Optical LEO",
            "Traffic-Aware Domain Partitioning & Inter-Domain Routing",
        ],
        "gap": "端到端可微流量工程空白;PINN 流量预测 + 守恒约束是 A 类潜力",
    },
    "06": {
        "core": "如何让全流水线对星座参数可微,梯度优化覆盖/时延?(最核心创新)",
        "benchmarks": "ESA ML-dSGP4 精度 +34% · Enzyme/Zygote AD · Adam lr=3e-3",
        "highlights": [
            "Orbital Prediction with Automatic Differentiable Perturbation Model",
            "Thermal Surrogate Model with Physics-Informed ML + POD",
            "Service-Differentiable Satellite Networking(IWQoS)",
            "Gradient-Based Spacecraft Thermal/Aircraft Design",
        ],
        "gap": "⚠️ 领域空白:可微 J2 + 软覆盖 loss + Adam 闭环无先例,本项目 optimize_coverage 属首创",
    },
    "07": {
        "core": "能否用 PINN 替代传统传播器,既匹配数据又满足运动方程?",
        "benchmarks": "Raissi 2019 PINN 奠基 · DeepXDE · FNO/DeepONet · 2024 首篇卫星 PINN",
        "highlights": [
            "Physics-Informed NN for Satellite State Estimation(2024 首篇)",
            "Tracking Space Debris Using PINN",
            "Learning Robust Satellite Attitude Dynamics(PINF)",
            "Physics-Informed ML for Real-Time Spacecraft Thermal",
        ],
        "gap": "PINN+卫星路由(0篇)、PINN+星上计算(0篇)、PINN+流量预测(0篇)全空白",
    },
    "08": {
        "core": "如何用 LLM/Agent 把自然语言翻译成仿真工具调用?",
        "benchmarks": "GPT-4/Claude 工具调用 · 'Language models are spacecraft operators'",
        "highlights": [
            "Language Models are Spacecraft Operators",
            "Generative AI Agents with LLM for Satellite Networks (MoE)",
            "Fine-tuning LLMs for Autonomous Spacecraft Control",
            "Leveraging LLMs for Integrated Satellite-Aerial-Terrestrial",
        ],
        "gap": "LLM 驱动的卫星仿真编排器(agent_repl)无直接对标",
    },
    "09": {
        "core": "如何度量切换频次/中断/乒乓率,设计低中断切换策略?",
        "benchmarks": "Hypatia 瞬时切换(简化) · 可见窗口 5-15min · OpenSN 自定义策略",
        "highlights": [
            "Trends in LEO Satellite Handover Algorithms",
            "Cooperative Beam Hopping for Positioning in Ultra-Dense LEO",
            "Proactive TCP Mechanism to Improve Handover Performance",
            "Joint User Association in Integrated Satellite-HAPS-Ground",
        ],
        "gap": "切换与路由耦合的中断度量较少;可微切换策略优化空白",
    },
    "10": {
        "core": "LEO 高 BDP 与移动性导致传统 TCP 性能下降,如何应对?",
        "benchmarks": "SaTCP(INFOCOM23) · LeoTCP(arXiv25) · BBR vs Cubic · Mathis 上界",
        "highlights": [
            "SaTCP: LEO 链路自适应 TCP(INFOCOM 2023)",
            "LeoTCP: LEO 专用 TCP(arXiv 2025)",
            "QPEP: QUIC-based Performance Enhancing Proxy",
            "Survey of Congestion Control in Satellite Networks",
        ],
        "gap": "本项目不自实现 TCP(保持简洁),用解析模型给上界 + ns-3 trace 接口",
    },
}


def slide_section_overview(prs, sm):
    """板块概览页(1页/板块)。"""
    sid = sm["id"]
    name = sm["name"]
    info = SECTION_OVERVIEW[sid]
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    color_hex = SECTION_COLORS[int(sid) - 1]
    bar_color = RGBColor(int(color_hex[1:3], 16), int(color_hex[3:5], 16),
                         int(color_hex[5:7], 16))
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(7.5), bar_color)
    # 板块编号大字
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(2), Inches(1.2),
                sid, font_size=64, bold=True, color=bar_color)
    add_textbox(slide, Inches(2.2), Inches(0.4), Inches(10), Inches(0.7),
                name, font_size=32, bold=True, color=COLOR_TEXT)
    add_textbox(slide, Inches(2.2), Inches(1.1), Inches(10), Inches(0.4),
                f"论文数:{sm['total']} 篇  "
                f"(★核心 {sm['tier1']} / ☆相关 {sm['tier2']} / ○借鉴 {sm['tier3']})",
                font_size=14, color=COLOR_MUTED)
    # 核心问题
    add_rect(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.9),
             RGBColor(0xEA, 0xF2, 0xFA))
    add_textbox(slide, Inches(0.7), Inches(1.85), Inches(2), Inches(0.4),
                "核心问题", font_size=13, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12), Inches(0.45),
                info["core"], font_size=15, color=COLOR_TEXT)
    # 对标基准
    add_textbox(slide, Inches(0.5), Inches(2.95), Inches(12), Inches(0.4),
                "▎对标基准", font_size=15, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.6), Inches(3.35), Inches(12.2), Inches(0.5),
                info["benchmarks"], font_size=13, color=COLOR_TEXT)
    # 代表论文
    add_textbox(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.4),
                "▎代表论文(亮点)", font_size=15, bold=True, color=COLOR_PRIMARY)
    for i, h in enumerate(info["highlights"]):
        y = 4.45 + i * 0.42
        add_rect(slide, Inches(0.6), Inches(y + 0.05), Inches(0.12), Inches(0.25),
                 bar_color)
        add_textbox(slide, Inches(0.85), Inches(y), Inches(12), Inches(0.4),
                    h, font_size=13, color=COLOR_TEXT)
    # 研究空白
    add_rect(slide, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.75),
             RGBColor(0xFF, 0xF4, 0xEA))
    add_textbox(slide, Inches(0.7), Inches(6.3), Inches(2), Inches(0.4),
                "研究空白", font_size=13, bold=True, color=COLOR_WARN)
    add_textbox(slide, Inches(0.7), Inches(6.62), Inches(12), Inches(0.35),
                info["gap"], font_size=12, color=COLOR_TEXT)


def slide_section_top_papers(prs, sm, items):
    """板块 Top 论文页(1页/板块)。"""
    sid = sm["id"]
    name = sm["name"]
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, f"板块 {sid} · {name} · 代表论文",
                  f"共 {sm['total']} 篇,以下为按相关性+年份排序的精选")
    # 表格:取前 8 篇
    top = items[:8]
    # 表头
    headers = ["#", "相关性", "年份", "来源", "标题"]
    col_w = [Inches(0.5), Inches(1.0), Inches(0.8), Inches(1.2), Inches(8.8)]
    x_start = Inches(0.5)
    y = Inches(1.5)
    row_h = Inches(0.55)
    # 表头行
    add_rect(slide, x_start, y, sum(col_w, Emu(0)), row_h, COLOR_PRIMARY)
    x = x_start
    for h, w in zip(headers, col_w):
        add_textbox(slide, x, y, w, row_h, h, font_size=12, bold=True,
                    color=COLOR_TEXT_LIGHT, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += w
    # 数据行
    tier_map = {"tier1": "★核心", "tier2": "☆相关", "tier3": "○借鉴"}
    tier_color = {"tier1": COLOR_DANGER, "tier2": COLOR_WARN,
                  "tier3": COLOR_MUTED}
    src_map = {"arXiv": "arXiv", "CCF": "CCF", "CCF_Conf": "CCF会议",
               "CAS": "CAS", "Aero": "航天", "IntlSup": "国际供", "Sim": "仿真"}
    import re
    for i, it in enumerate(top):
        ry = y + row_h * (i + 1)
        bg = RGBColor(0xF5, 0xF7, 0xFA) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        add_rect(slide, x_start, ry, sum(col_w, Emu(0)), row_h, bg,
                 RGBColor(0xDD, 0xDD, 0xDD))
        x = x_start
        # 序号
        add_textbox(slide, x, ry, col_w[0], row_h, str(i + 1), font_size=11,
                    color=COLOR_MUTED, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[0]
        # 相关性
        t = tier_map.get(it["tier"], it["tier"])
        add_textbox(slide, x, ry, col_w[1], row_h, t, font_size=11, bold=True,
                    color=tier_color.get(it["tier"], COLOR_TEXT),
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[1]
        # 年份
        add_textbox(slide, x, ry, col_w[2], row_h, it["year"] or "-",
                    font_size=11, color=COLOR_TEXT, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[2]
        # 来源
        add_textbox(slide, x, ry, col_w[3], row_h,
                    src_map.get(it["source"], it["source"]),
                    font_size=11, color=COLOR_TEXT, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[3]
        # 标题
        title = re.sub(r"^\s*\d+\.\s*", "", it["title"]).strip()
        add_textbox(slide, x + Pt(4), ry, col_w[4] - Pt(8), row_h, title,
                    font_size=11, color=COLOR_TEXT, anchor=MSO_ANCHOR.MIDDLE)
    # 提示
    add_textbox(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
                f"💡 完整 {sm['total']} 篇论文清单见 docs/literature/ 目录下对应板块 .md 文件",
                font_size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)


def slide_gaps(prs):
    """三大研究空白页。"""
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, "05  研究空白与本项目创新",
                  "三大蓝海方向 × SatelliteSimJulia 对应实现")
    gaps = [
        ("🥇", "可微仿真闭环", "板块06",
         "可微 J2/SGP4 → 软 ISL/覆盖 loss → Adam 端到端星座优化",
         "严格可微论文仅 17 篇,卫星领域几乎空白",
         "src/opt 的 optimize_coverage driver + Enzyme/Zygote",
         COLOR_DANGER),
        ("🥈", "PINN + 卫星网络", "板块07",
         "PINN 路由时延预测器 / PINN 星上计算 / PINN 流量预测",
         "PINN+卫星路由/星上计算/流量预测均为 0 篇",
         "src/opt 的 pinn_routing.jl + pinn_model.jl",
         COLOR_WARN),
        ("🥉", "LLM 仿真编排", "板块08",
         "LLM 把自然语言翻译成仿真工具调用,意图防泄漏",
         "LLM + 卫星仿真工具链集成无成熟编排器",
         "src/lab 的 SimAgent + agent_repl + Intent 翻译",
         COLOR_ACCENT),
    ]
    y = 1.5
    for medal, title, sec, desc, gap, impl, color in gaps:
        # 左侧色条
        add_rect(slide, Inches(0.5), Inches(y), Inches(0.15), Inches(1.65), color)
        # 奖牌
        add_textbox(slide, Inches(0.75), Inches(y), Inches(0.7), Inches(0.7),
                    medal, font_size=36, anchor=MSO_ANCHOR.MIDDLE)
        # 标题
        add_textbox(slide, Inches(1.5), Inches(y), Inches(4), Inches(0.5),
                    title, font_size=20, bold=True, color=COLOR_TEXT)
        add_textbox(slide, Inches(5.6), Inches(y + 0.05), Inches(1.5), Inches(0.4),
                    sec, font_size=13, bold=True, color=color,
                    align=PP_ALIGN.CENTER)
        # 描述
        add_textbox(slide, Inches(1.5), Inches(y + 0.55), Inches(11), Inches(0.4),
                    f"方向:{desc}", font_size=13, color=COLOR_TEXT)
        # 空白
        add_textbox(slide, Inches(1.5), Inches(y + 0.95), Inches(11), Inches(0.4),
                    f"空白:{gap}", font_size=13, color=COLOR_MUTED)
        # 实现
        add_textbox(slide, Inches(1.5), Inches(y + 1.3), Inches(11), Inches(0.35),
                    f"实现:{impl}", font_size=12, bold=True, color=COLOR_PRIMARY)
        y += 1.85


def slide_roadmap(prs):
    """路线图页。"""
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, "06  推荐研究路线图",
                  "基于文献饱和度与空白识别,分阶段推进")
    # 三阶段
    stages = [
        ("阶段 1 · 短期", "对标验证", COLOR_PRIMARY, [
            "复现 Hypatia 基准(RTT/跳数/MLU) → 板块04,05",
            "传播器精度对比(TwoBody/J2/SGP4) → 板块01",
            "ISL/拓扑评估验证 → 板块02,03",
        ]),
        ("阶段 2 · 中期", "核心创新", COLOR_WARN, [
            "★ 可微 J2 + 软覆盖 loss + Adam 闭环 → 板块06",
            "　　 (第一篇论文)",
            "切换中断度量建模 → 板块09",
            "流量工程可微化探索 → 板块05",
        ]),
        ("阶段 3 · 长期", "蓝海突破", COLOR_DANGER, [
            "★★ PINN 路由时延预测器 → 板块07 (第二篇论文)",
            "★★ PINN + 星上计算/流量预测 → 板块07 (第三篇)",
            "LLM 仿真编排器产品化 → 板块08",
        ]),
    ]
    col_w = Inches(4.0)
    gap = Inches(0.2)
    x = Inches(0.5)
    y = Inches(1.5)
    for title, sub, color, bullets in stages:
        # 卡片背景
        add_rect(slide, x, y, col_w, Inches(5.2),
                 RGBColor(0xF5, 0xF7, 0xFA), color)
        # 顶部色块
        add_rect(slide, x, y, col_w, Inches(0.9), color)
        add_textbox(slide, x, y + Inches(0.1), col_w, Inches(0.5),
                    title, font_size=20, bold=True, color=COLOR_TEXT_LIGHT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.55), col_w, Inches(0.35),
                    sub, font_size=14, color=COLOR_TEXT_LIGHT,
                    align=PP_ALIGN.CENTER)
        # 内容
        add_bullets(slide, x + Inches(0.2), y + Inches(1.15),
                    col_w - Inches(0.4), Inches(3.8), bullets,
                    font_size=13, line_spacing=1.4)
        x += col_w + gap


def slide_summary(prs, data):
    """总结页。"""
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), COLOR_ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
                "总结与下一步", font_size=40, bold=True, color=COLOR_TEXT_LIGHT)
    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
                "SatelliteSimJulia 文献调研全景汇报", font_size=18,
                color=COLOR_ACCENT)
    # 核心结论
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.5),
                "▎核心结论", font_size=20, bold=True, color=COLOR_ACCENT)
    add_bullets(slide, Inches(0.9), Inches(2.75), Inches(11.5), Inches(2.5), [
        f"从 {data['meta']['total_papers_in_db']:,} 篇论文库筛选出 "
        f"{data['meta']['total_unique_matched']:,} 篇相关文献,按 10 个技术层级组织",
        "路由/ISL/切换(板块02/04/09)是成熟红海,可对标 Hypatia/OpenSN 直接复现",
        "可微优化(板块06)、PINN(板块07)、LLM编排(板块08)是三大蓝海,具备 A 类论文潜力",
        "本项目 src/opt(pinn + 可微优化)与 src/lab(LLM agent)正处于空白地带,先发优势明显",
    ], font_size=15, color=COLOR_TEXT_LIGHT, line_spacing=1.5)
    # 下一步
    add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
                "▎下一步", font_size=20, bold=True, color=COLOR_ACCENT)
    add_bullets(slide, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.5), [
        "详细文献清单见 docs/literature/ 各板块 .md(全量论文 + 子主题分组)",
        "优先完成可微优化闭环(阶段2)→ 出第一篇论文",
        "PINN 方向作为第二、三篇论文储备",
    ], font_size=14, color=COLOR_TEXT_LIGHT, line_spacing=1.4)
    # 底部
    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                "SatelliteSimJulia  ·  LEO 卫星星座仿真 + 可微优化 + AI 适配",
                font_size=12, color=COLOR_MUTED)


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------
def main():
    if not os.path.exists(JSON_IN):
        sys.exit(f"JSON 不存在:{JSON_IN},请先运行 build_literature_index.py")

    with open(JSON_IN, "r", encoding="utf-8") as f:
        data = json.load(f)

    secs_meta = data["sections_meta"]
    secs = data["sections"]

    # 预生成图表
    print("生成图表...")
    chart_count = make_section_count_chart(secs_meta)
    chart_tier = make_tier_stacked_chart(secs_meta)
    chart_year = make_year_trend_chart(data)
    chart_heat = make_heatmap_chart(secs_meta)

    # 创建 PPT (16:9)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("构建幻灯片...")
    # 1. 封面
    slide_cover(prs, data)
    # 2. 目录
    slide_toc(prs)
    # 3. 背景
    slide_background(prs)
    # 4. 方法论
    slide_methodology(prs, data)
    # 5-6. 全景统计
    slide_stats_overview(prs, data, chart_count, chart_tier)
    slide_year_trend(prs, data, chart_year)
    # 7. 热度图(独立一页)
    slide = add_slide(prs)
    set_bg(slide, COLOR_BG_LIGHT)
    add_title_bar(slide, "04  十大板块文献全景",
                  "研究热度排序:红海(饱和)→ 绿海(蓝海)")
    slide.shapes.add_picture(chart_heat, Inches(2), Inches(1.4),
                             width=Inches(9.3))
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
                "图例:红色 ≥400篇(高度饱和) · 橙色 150-400(活跃) · 绿色 <150(蓝海/新兴)",
                font_size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
    # 8-27. 10 板块各 2 页
    for sm in secs_meta:
        items = secs.get(sm["id"], [])
        slide_section_overview(prs, sm)
        slide_section_top_papers(prs, sm, items)
    # 研究空白
    slide_gaps(prs)
    # 路线图
    slide_roadmap(prs)
    # 总结
    slide_summary(prs, data)

    # 保存
    prs.save(PPT_OUT)
    print(f"\n✓ PPT 已生成:{PPT_OUT}")
    print(f"  共 {len(prs.slides)} 页")


if __name__ == "__main__":
    main()
